import base64
import io
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation

def extract_text_from_base64(file_data_b64: str, file_type: str) -> str:
    # Remove data URI header if present (e.g. "data:application/pdf;base64,...")
    if "," in file_data_b64:
        file_data_b64 = file_data_b64.split(",", 1)[1]
        
    try:
        file_bytes = base64.b64decode(file_data_b64)
    except Exception as e:
        return f"[Error decoding file: {e}]"
        
    file_stream = io.BytesIO(file_bytes)
    
    text = ""
    try:
        if file_type == "pdf":
            reader = PdfReader(file_stream)
            for page in reader.pages:
                extracted = page.extract_text()
                if extracted:
                    text += extracted + "\n"
        elif file_type in ["docx", "doc"]:
            doc = Document(file_stream)
            for para in doc.paragraphs:
                text += para.text + "\n"
        elif file_type == "pptx":
            prs = Presentation(file_stream)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        elif file_type == "txt":
            text = file_bytes.decode('utf-8')
        else:
            return f"[Unsupported document format: {file_type}]"
    except Exception as e:
        return f"[Error parsing document: {e}]"
        
    return text.strip()
