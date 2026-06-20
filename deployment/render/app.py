import os
import time
import base64
import io
import diskcache
import uuid
from typing import Optional, List, Dict
from pydantic import BaseModel
from fastapi import FastAPI, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from openai import OpenAI
from dotenv import load_dotenv
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation

load_dotenv()

# ==========================================
# 1. DOCUMENT PARSER
# ==========================================
def extract_text_from_base64(file_data_b64: str, file_type: str) -> str:
    if "," in file_data_b64:
        file_data_b64 = file_data_b64.split(",", 1)[1]
    try:
        file_bytes = base64.b64decode(file_data_b64)
    except Exception as e:
        return f"[Error decoding file: {e}]"
    file_stream = io.BytesIO(file_bytes)
    text = ""
    try:
        if file_type == "pdf":
            reader = PdfReader(file_stream)
            for page in reader.pages:
                extracted = page.extract_text()
                if extracted:
                    text += extracted + "\n"
        elif file_type in ["docx", "doc"]:
            doc = Document(file_stream)
            for para in doc.paragraphs:
                text += para.text + "\n"
        elif file_type == "pptx":
            prs = Presentation(file_stream)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        elif file_type == "txt":
            text = file_bytes.decode('utf-8')
        else:
            return f"[Unsupported document format: {file_type}]"
    except Exception as e:
        return f"[Error parsing document: {e}]"
    return text.strip()


# ==========================================
# 2. SOCRATIC LAYER (With Serverless Filter)
# ==========================================
class SessionManager:
    def __init__(self, cache_dir: str = './session_cache'):
        self.cache = diskcache.Cache(cache_dir)
        self.SESSION_TTL = 60 * 60 * 24

    def create_session(self, session_id: str, topic: str):
        data = {'history': [], 'stage': 1, 'subject': topic, 'hints_given': 0, 'created': time.time()}
        self.cache.set(session_id, data, expire=self.SESSION_TTL)
        return data

    def get_session(self, session_id: str) -> Optional[dict]:
        return self.cache.get(session_id, default=None)

    def update_session(self, session_id: str, data: dict):
        self.cache.set(session_id, data, expire=self.SESSION_TTL)

PROMPT_TEMPLATES = {
    1: 'You are a Socratic tutor. NEVER reveal the answer if the student does not know it. If they answer correctly, warmly confirm it and move to the next concept. Otherwise, ask a broad exploratory question. Subject: {subject}.',
    2: 'You are a Socratic tutor. NEVER reveal the answer. If the student answers correctly, confirm it. Otherwise, ask for their initial hypothesis. Subject: {subject}.',
    3: 'You are a Socratic tutor. NEVER reveal the answer. If they answer correctly, confirm it. Otherwise, probe for misconceptions with a guiding question. Subject: {subject}.',
    4: 'You are a Socratic tutor teaching {subject}. The student needs a hint. Provide a structural framework hint without giving the answer. End with a question.',
    5: 'You are a Socratic tutor teaching {subject}. The student is very close. Give a near-direct hint, but require them to make the final connection. End with a question.',
    6: 'The student appears to understand {subject}. Ask them to explain the concept in their own words. Do NOT provide a model answer.',
    7: 'The student has successfully understood {subject}. Affirm their understanding and ask if they have any other questions.'
}

class ConfidenceScorer:
    def score(self, student_response: str) -> float:
        score = 0.5
        if len(student_response.split()) < 5: score -= 0.2
        if len(student_response.split()) > 20: score += 0.2
        for marker in ['because', 'therefore', 'so', 'means']:
            if marker in student_response.lower(): score += 0.1
        if '?' in student_response: score -= 0.1
        return max(0.0, min(1.0, score))

class HintProgressionEngine:
    def __init__(self):
        self.scorer = ConfidenceScorer()

    def determine_next_stage(self, current_stage, student_response):
        confidence = self.scorer.score(student_response)
        if current_stage == 1: return 2
        elif current_stage == 2: return 3 if confidence < 0.6 else 6
        elif current_stage == 3:
            if confidence < 0.4: return 4
            if confidence >= 0.6: return 6
            return 3
        elif current_stage == 4:
            if confidence < 0.5: return 5
            if confidence >= 0.6: return 6
            return 4
        elif current_stage == 5: return 6 if confidence >= 0.7 else 5
        elif current_stage == 6: return 7 if confidence >= 0.7 else 2
        return current_stage

# ==========================================
# 3. FASTAPI SETUP
# ==========================================
app = FastAPI()

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=False,
    allow_methods=["*"],
    allow_headers=["*"],
)

session_manager = SessionManager()
hint_engine = HintProgressionEngine()

class NewSessionRequest(BaseModel):
    student_id: str
    topic: str

class ChatRequest(BaseModel):
    session_id: str
    student_id: str
    message: str
    file_data: Optional[str] = None
    file_type: Optional[str] = None
    system_prompt: Optional[str] = None

# ==========================================
# 4. HUGGING FACE SERVERLESS INFERENCE & FILTER
# ==========================================
# Make sure HF_TOKEN is set in your Render Environment Variables!
# You can generate a free token at https://huggingface.co/settings/tokens
hf_client = OpenAI(
    api_key=os.environ.get("HF_TOKEN", "missing_key"),
    base_url="https://api-inference.huggingface.co/v1/",
)
MODEL_ID = "Qwen/Qwen2.5-7B-Instruct"

def generate_response(messages: List[Dict]) -> str:
    text_messages = []
    for msg in messages:
        if isinstance(msg['content'], list):
            text_content = " ".join([c.get("text", "") for c in msg['content'] if c.get("type") == "text"])
            text_messages.append({"role": msg["role"], "content": text_content})
        else:
            text_messages.append({"role": msg["role"], "content": msg["content"]})
            
    try:
        response = hf_client.chat.completions.create(
            model=MODEL_ID,
            messages=text_messages,
            max_tokens=512,
            temperature=0.7
        )
        return response.choices[0].message.content
    except Exception as e:
        return f"[Error from Hugging Face API]: {str(e)}"

class ServerlessSocraticFilter:
    def blocks_direct_answer(self, ai_response: str) -> bool:
        prompt = f"""
        Analyze the following response from a tutor to a student.
        Did the tutor directly reveal the correct answer to a problem the student is trying to solve, instead of guiding them?
        Answer with ONLY 'Yes' or 'No'.
        
        Tutor Response:
        "{ai_response}"
        """
        try:
            response = hf_client.chat.completions.create(
                model=MODEL_ID,
                messages=[{"role": "user", "content": prompt}],
                max_tokens=5,
                temperature=0.1
            )
            answer = response.choices[0].message.content.strip().lower()
            return "yes" in answer
        except Exception:
            return False

socratic_filter = ServerlessSocraticFilter()

# ==========================================
# 5. FASTAPI ROUTES
# ==========================================
@app.post("/v1/session/new")
def new_session(req: NewSessionRequest):
    session_id = str(uuid.uuid4())
    session_manager.create_session(session_id, req.topic)
    return {"session_id": session_id}

@app.post("/v1/chat")
def chat(request: ChatRequest):
    session = session_manager.get_session(request.session_id)
    if not session:
        session = session_manager.create_session(request.session_id, "general")

    stage = session['stage']
    subject = session.get('subject', 'general')
    
    if request.system_prompt:
        system_content = request.system_prompt
    else:
        template = PROMPT_TEMPLATES.get(stage, PROMPT_TEMPLATES[1])
        system_content = template.format(subject=subject)
 
    messages = [{'role': 'system', 'content': system_content}]
    messages += session['history']
    
    user_message_content = request.message
    is_vision = False
    
    if request.file_data and request.file_type:
        ext = request.file_type.lower()
        if ext in ['pdf', 'docx', 'doc', 'pptx', 'txt']:
            doc_text = extract_text_from_base64(request.file_data, ext)
            if doc_text:
                user_message_content += f"\n\n[Attached Document Content]:\n{doc_text}"
        elif ext in ['png', 'jpg', 'jpeg', 'webp', 'gif']:
            is_vision = True
            user_content = [
                {"type": "text", "text": user_message_content},
                {"type": "image_url", "image_url": {"url": request.file_data}}
            ]
            messages.append({'role': 'user', 'content': user_content})

    if not is_vision:
        messages.append({'role': 'user', 'content': user_message_content})
 
    response = generate_response(messages)
 
    if socratic_filter.blocks_direct_answer(response):
        print("[FILTER] Model leaked answer! Intervening and requesting guiding question...")
        messages.append({'role': 'assistant', 'content': response})
        messages.append({'role': 'user', 'content': 'That was too direct. Please try again and ask a guiding question instead without revealing the answer.'})
        response = generate_response(messages)
 
    if is_vision:
        session['history'].append({'role': 'user', 'content': user_content})
    else:
        session['history'].append({'role': 'user', 'content': user_message_content})
    session['history'].append({'role': 'assistant', 'content': response})
    
    next_stage = hint_engine.determine_next_stage(stage, request.message)
    session['stage'] = next_stage
    
    if next_stage in [3, 4, 5]:
        session['hints_given'] += 1

    session_manager.update_session(request.session_id, session)
    confidence = hint_engine.scorer.score(request.message)

    return {
        "response": response, 
        "stage": next_stage,
        "hints_given": session['hints_given'],
        "confidence": confidence
    }

@app.get("/v1/session/{session_id}")
def get_session(session_id: str):
    session = session_manager.get_session(session_id)
    if not session:
        raise HTTPException(status_code=404, detail="Session not found")
    return session

@app.get("/health")
def health():
    return {"status": "ok", "backend": "Render+TogetherAI"}

if __name__ == "__main__":
    import uvicorn
    # Render binds the port via the PORT environment variable
    port = int(os.environ.get("PORT", 8000))
    uvicorn.run(app, host="0.0.0.0", port=port)
